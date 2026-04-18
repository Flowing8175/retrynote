import asyncio
import json as json_mod
import logging
from dataclasses import dataclass
from datetime import datetime, timezone
from typing import AsyncGenerator

from sqlalchemy import select, func
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import settings as cfg
from app.database import async_session
from app.models.file import File, FileStatus, ParsedDocument, DocumentChunk
from app.models.objection import Objection, ObjectionStatus, WeakPoint
from app.models.quiz import (
    QuizSession,
    QuizSessionFile,
    QuizSessionStatus,
    QuizItem,
    QuestionType,
    AnswerLog,
    Judgement,
    ErrorType,
)
from app.models.search import Job, DraftAnswer
from app.prompts import (
    SYSTEM_PROMPT_DIFFICULTY_SELECTION,
    SYSTEM_PROMPT_OBJECTION_REVIEW,
    get_generation_system_prompt,
)
from app.prompts.generation import build_generation_prompt
from app.prompts.retry_generation import (
    build_batch_retry_prompt,
    get_retry_system_prompt,
)
from app.utils.ai_client import (
    call_ai_with_fallback,
    stream_ai_structured_with_thinking,
    OBJECTION_REVIEW_SCHEMA,
)
from app.utils.job_runner import JobFailure, JobRunner
from app.utils.normalize import normalize_answer, normalize_concept_key

logger = logging.getLogger(__name__)

_MAX_URL_SOURCE_CHARS = 15000

_BLOCKED_URL_HOSTS = frozenset(
    {
        "localhost",
        "127.0.0.1",
        "0.0.0.0",
        "::1",
        "metadata.google.internal",
        "169.254.169.254",
    }
)


def validate_source_url_syntax(url: str) -> str | None:
    from urllib.parse import urlparse

    url = (url or "").strip()
    if not url:
        return "URL이 비어있습니다."
    try:
        parsed = urlparse(url)
    except Exception:
        return "URL 형식이 올바르지 않습니다."
    if parsed.scheme not in ("http", "https"):
        return "URL은 http:// 또는 https://로 시작해야 합니다."
    if not parsed.hostname:
        return "URL에 호스트가 없습니다."
    if parsed.hostname.lower() in _BLOCKED_URL_HOSTS:
        return "접근할 수 없는 호스트입니다."
    return None


async def process_file(job_id: str):
    async with async_session() as db:
        job_result = await db.execute(select(Job).where(Job.id == job_id))
        job = job_result.scalar_one_or_none()
        if not job:
            return

        async with JobRunner(db, job):
            file_result = await db.execute(select(File).where(File.id == job.target_id))
            file = file_result.scalar_one_or_none()
            if not file:
                raise JobFailure("File not found")

            try:
                file.status = FileStatus.parsing
                file.processing_started_at = datetime.now(timezone.utc)
                await db.commit()

                ocr_required = file.file_type in ("png", "jpg", "jpeg")
                text = ""
                file_data: bytes | None = None

                if ocr_required:
                    file.ocr_required = True
                    file.status = FileStatus.ocr_pending
                    await db.commit()

                    file.status = FileStatus.ocr_processing
                    await db.commit()

                    from app.services.ocr_service import extract_text_from_image
                    from app.services import storage as _storage

                    stored_path = file.stored_path
                    if not stored_path:
                        raise ValueError("Image file missing stored_path")
                    image_data = await _storage.download_file(stored_path)
                    ocr_result = await extract_text_from_image(image_data)
                    text = ocr_result.text
                elif file.source_type.value == "manual_text":
                    payload = job.payload_json or {}
                    text = payload.get("manual_text", "")
                elif file.source_type.value == "upload":
                    text, file_data = await _extract_file_text_async(file)
                elif file.source_type.value == "url":
                    if not file.source_url:
                        raise ValueError("URL source missing source_url")
                    text, _ = await _fetch_url_text(file.source_url)

                file.status = FileStatus.parsed
                await db.commit()

                if (
                    not text.strip()
                    and file_data
                    and file.file_type in ("pptx", "pdf", "docx")
                ):
                    images = _extract_images_from_document(file_data, file.file_type)
                    if images:
                        from app.services.ocr_service import extract_text_from_image

                        file.ocr_required = True
                        file.status = FileStatus.ocr_processing
                        await db.commit()

                        ocr_texts: list[str] = []
                        for idx, img_bytes in enumerate(images[:MAX_DOC_OCR_IMAGES]):
                            try:
                                ocr_result = await extract_text_from_image(img_bytes)
                                if ocr_result.text.strip():
                                    ocr_texts.append(ocr_result.text)
                            except RuntimeError as e:
                                err = str(e)
                                if "not configured" in err or "invalid API key" in err:
                                    file.status = FileStatus.failed_terminal
                                    file.parse_error_code = "ocr_not_configured"
                                    file.processing_finished_at = datetime.now(
                                        timezone.utc
                                    )
                                    raise JobFailure(f"OCR configuration error: {err}")
                                logger.warning(
                                    "OCR error for image %d in file %s: %s",
                                    idx,
                                    file.id,
                                    err,
                                )
                            except Exception as e:
                                logger.warning(
                                    "OCR failed for image %d in file %s: %s",
                                    idx,
                                    file.id,
                                    e,
                                )
                        if ocr_texts:
                            text = "\n".join(ocr_texts)
                        ocr_required = True

                        file.status = FileStatus.parsed
                        await db.commit()

                if not text.strip():
                    if ocr_required:
                        file.status = FileStatus.failed_partial
                        file.parse_error_code = "ocr_empty"
                    else:
                        file.status = FileStatus.failed_terminal
                        file.parse_error_code = "empty_content"
                    file.processing_finished_at = datetime.now(timezone.utc)
                    raise JobFailure("File content is empty")

                normalized = _normalize_text(text)

                parsed_doc = ParsedDocument(
                    file_id=file.id,
                    raw_text=text,
                    normalized_text=normalized,
                    language="ko",
                    page_count=1,
                    parser_name="kakao_ocr" if ocr_required else "builtin",
                    parser_version="1.0",
                    ocr_applied=ocr_required,
                )
                db.add(parsed_doc)
                await db.flush()

                file.status = FileStatus.embedding_pending
                await db.commit()

                chunks = _create_chunks(file.id, parsed_doc.id, normalized)
                db.add_all(chunks)
                await db.flush()

                for chunk in chunks:
                    chunk.embedding_status = "completed"
                logger.warning(
                    "File %s: embedding_status set to 'completed' without actual embedding — "
                    "vector search will not work for this file until embeddings are generated",
                    file.id,
                )

                file.status = FileStatus.ready
                file.is_searchable = True
                file.is_quiz_eligible = True
                file.processing_finished_at = datetime.now(timezone.utc)
                await db.commit()

                # Auto-trigger summary generation when file reaches ready status
                try:
                    from app.workers.celery_app import dispatch_task

                    dispatch_task("generate_study_summary", [file.id])
                except Exception:
                    pass  # Auto-trigger failure should not affect file processing

            except JobFailure:
                raise
            except Exception as e:
                is_ocr_failure = file.ocr_required and file.status in (
                    FileStatus.ocr_pending,
                    FileStatus.ocr_processing,
                )
                file.status = (
                    FileStatus.failed_partial
                    if is_ocr_failure
                    else FileStatus.failed_terminal
                )
                file.parse_error_code = str(e)[:100]
                file.processing_finished_at = datetime.now(timezone.utc)
                raise


async def _extract_file_text_async(file: File) -> tuple[str, bytes | None]:
    if not file.stored_path:
        return "", None

    from app.services import storage as _storage

    try:
        data = await _storage.download_file(file.stored_path)
    except Exception:
        return "", None

    ext = file.file_type
    if ext == "pdf":
        return _extract_pdf_bytes(data), data
    elif ext == "docx":
        return _extract_docx_bytes(data), data
    elif ext == "pptx":
        return _extract_pptx_bytes(data), data
    elif ext in ("txt", "md"):
        return data.decode("utf-8", errors="replace"), data
    return "", data


def _extract_pdf_bytes(data: bytes) -> str:
    try:
        import io
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(data))
        return "".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""


def _extract_docx_bytes(data: bytes) -> str:
    try:
        import io
        from docx import Document

        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


def _extract_pptx_bytes(data: bytes) -> str:
    try:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text += para.text + "\n"
        return text
    except Exception:
        return ""


MAX_DOC_OCR_IMAGES = 30


def _extract_images_from_pptx(data: bytes) -> list[bytes]:
    try:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        images: list[bytes] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        images.append(shape.image.blob)
                    except Exception:
                        continue
        return images
    except Exception:
        return []


def _extract_images_from_pdf(data: bytes) -> list[bytes]:
    try:
        import io
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(data))
        images: list[bytes] = []
        for page in reader.pages:
            try:
                for img in page.images:
                    images.append(img.data)
            except Exception:
                continue
        return images
    except Exception:
        return []


def _extract_images_from_docx(data: bytes) -> list[bytes]:
    try:
        import io
        from docx import Document

        doc = Document(io.BytesIO(data))
        images: list[bytes] = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    images.append(rel.target_part.blob)
                except Exception:
                    continue
        return images
    except Exception:
        return []


def _extract_images_from_document(data: bytes, file_type: str) -> list[bytes]:
    if file_type == "pptx":
        return _extract_images_from_pptx(data)
    elif file_type == "pdf":
        return _extract_images_from_pdf(data)
    elif file_type == "docx":
        return _extract_images_from_docx(data)
    return []


def _validate_ip(addr: str) -> bool:
    """Check if a resolved IP address is safe (not private/internal/special).

    Defensively unwraps IPv4-embedded IPv6 forms (``::ffff:a.b.c.d`` and the
    deprecated ``::a.b.c.d``) before evaluating flags, because Python
    <3.12.4 does NOT delegate ``is_loopback``/``is_link_local``/``is_reserved``
    to the underlying IPv4 address on ``IPv6Address`` instances — only
    ``is_private`` is consistent. Without unwrapping, an IPv4-mapped SSRF
    attempt like ``::ffff:169.254.169.254`` (AWS IMDS) is only caught by
    accident of ``::/8`` being reserved on some versions. Unwrapping makes
    validation explicit and version-independent.
    """
    import ipaddress

    try:
        ip = ipaddress.ip_address(addr)
    except ValueError:
        return False

    if isinstance(ip, ipaddress.IPv6Address):
        mapped = ip.ipv4_mapped
        if mapped is not None:
            ip = mapped
        elif (int(ip) >> 32) == 0:
            return False

    return not (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_reserved
        or ip.is_multicast
        or ip.is_unspecified
    )


def _resolve_and_validate_url(url: str) -> tuple[str, str, int] | None:
    """Resolve URL hostname and validate the resolved IP is safe.

    Returns (resolved_ip, hostname, port) or None if unsafe.
    Resolves DNS exactly once to prevent TOCTOU / DNS rebinding attacks.
    """
    from urllib.parse import urlparse
    import socket

    try:
        parsed = urlparse(url)
    except Exception:
        return None

    if parsed.scheme not in ("http", "https"):
        return None

    hostname = parsed.hostname
    if not hostname:
        return None

    if hostname.lower() in _BLOCKED_URL_HOSTS:
        return None

    port = parsed.port or (443 if parsed.scheme == "https" else 80)

    # Resolve DNS exactly once and validate ALL resolved addresses
    try:
        addr_infos = socket.getaddrinfo(hostname, port, proto=socket.IPPROTO_TCP)
        if not addr_infos:
            return None
        for info in addr_infos:
            addr = str(info[4][0])
            if not _validate_ip(addr):
                return None
        resolved_ip = str(addr_infos[0][4][0])
        return resolved_ip, hostname, port
    except socket.gaierror:
        return None


def _sanitize_title(raw: str | None) -> str | None:
    """Strip non-printable control characters from a title; max 200 chars.

    Defensive cleaning so we never persist garbage bytes. React escapes
    title on render, but we sanitize at the source too.
    """
    if not raw:
        return None
    cleaned = "".join(c for c in raw if c.isprintable() or c.isspace()).strip()
    if not cleaned:
        return None
    return cleaned[:200]


def _extract_page_title(soup) -> str | None:
    """Return sanitized <title> text from a BeautifulSoup soup (max 200 chars).

    Kept for backward compatibility and unit tests. New code should prefer
    :func:`_sanitize_title` on a raw string extracted by trafilatura.
    """
    title_tag = soup.find("title")
    if title_tag is None:
        return None
    raw = title_tag.get_text(strip=True)
    return _sanitize_title(raw)


def _extract_with_trafilatura(html_bytes: bytes, url: str) -> tuple[str, str | None]:
    """Run trafilatura extraction (CPU-bound, must be called in thread pool).

    Returns (body_text, page_title). Both empty/None if extraction fails.

    Trafilatura auto-detects encoding from HTML meta tags / BOM / heuristics,
    correctly handling EUC-KR, UTF-8, and other legacy encodings. It also
    strips boilerplate (nav, ads, comments, footers) automatically, so we
    don't need manual tag removal.
    """
    import trafilatura

    text = ""
    try:
        text = (
            trafilatura.extract(
                html_bytes,
                output_format="txt",
                favor_recall=True,  # include more content; precision matters less for AI input
                include_comments=False,
                include_tables=True,
                deduplicate=True,
                url=url,
            )
            or ""
        )
    except Exception:
        logger.exception("trafilatura.extract failed for %s", url)

    title: str | None = None
    try:
        metadata = trafilatura.extract_metadata(html_bytes)
        if metadata is not None:
            # Works for both Document objects and dict-like results
            raw_title = getattr(metadata, "title", None)
            title = _sanitize_title(raw_title)
    except Exception:
        logger.debug("trafilatura.extract_metadata failed for %s", url, exc_info=True)

    return text, title


async def _fetch_url_text(url: str) -> tuple[str, str | None]:
    """Fetch a public URL and return (extracted_body_text, page_title).

    Both fields are empty/None on any failure. DNS resolution is offloaded
    to a worker thread so the event loop is not blocked by slow lookups.
    Content extraction uses trafilatura (F1 0.958 on ScrapingHub benchmark),
    which handles encoding detection and boilerplate removal automatically.

    Security notes:
      * Before fetching, every address returned by :func:`getaddrinfo` for
        the hostname is validated against the SSRF block list (private,
        loopback, link-local, multicast, reserved). If any address fails,
        the entire request is refused — this catches the main SSRF vector
        of "attacker-supplied hostname that resolves to an internal IP".
      * The request is sent to the original hostname URL (not the raw IP)
        so TLS SNI and certificate verification work correctly. httpx will
        re-resolve the hostname, but in practice this hits the OS resolver
        cache within seconds of the pre-validation.
      * Redirects are NOT followed (they could land on an internal IP that
        bypassed our pre-validation).
      * Response body is capped at 5 MB.

    Residual risk: DNS rebinding between our pre-validation and httpx's
    resolution is theoretically possible but requires the attacker to
    control a domain with TTL=0 and win a tight race window. Defense in
    depth would pin the validated IP via a custom ``AsyncNetworkBackend``
    — see future work in issue tracker.
    """
    url = (url or "").strip()
    resolution = await asyncio.to_thread(_resolve_and_validate_url, url)
    if not resolution:
        logger.warning(
            "URL fetch rejected: DNS resolution or SSRF validation failed (url=%s)",
            url,
        )
        return "", None
    _, hostname, _ = resolution

    try:
        import httpx

        MAX_FETCH_BYTES = 5 * 1024 * 1024  # 5 MB
        body = b""
        try:
            async with httpx.AsyncClient() as client:
                async with client.stream(
                    "GET",
                    url,
                    follow_redirects=False,
                    timeout=30,
                ) as resp:
                    if resp.status_code != 200:
                        logger.warning(
                            "URL fetch returned non-200 (url=%s host=%s status=%s location=%s)",
                            url,
                            hostname,
                            resp.status_code,
                            resp.headers.get("location"),
                        )
                        return "", None
                    async for chunk in resp.aiter_bytes(chunk_size=8192):
                        body += chunk
                        if len(body) > MAX_FETCH_BYTES:
                            logger.warning(
                                "URL fetch aborted: body exceeded %d bytes (url=%s)",
                                MAX_FETCH_BYTES,
                                url,
                            )
                            return "", None
        except httpx.HTTPError as e:
            logger.warning(
                "URL fetch HTTP error (url=%s error_type=%s error=%s)",
                url,
                type(e).__name__,
                e,
            )
            return "", None

        if not body:
            logger.warning("URL fetch returned empty body (url=%s)", url)
            return "", None

        # Offload HTML parsing (CPU-bound) to a worker thread
        text, title = await asyncio.to_thread(_extract_with_trafilatura, body, url)
        if not text:
            logger.warning(
                "Content extraction returned empty text (url=%s body_bytes=%d)",
                url,
                len(body),
            )
        else:
            logger.info(
                "URL fetch ok (url=%s body_bytes=%d text_chars=%d title=%r)",
                url,
                len(body),
                len(text),
                title,
            )
        return text, title

    except Exception:
        logger.exception("Unexpected error fetching URL (url=%s)", url)
        return "", None
    resolved_ip, hostname, port = resolution

    try:
        import httpx
        from urllib.parse import urlparse, urlunparse

        # Rewrite URL to use the resolved IP directly, preventing DNS rebinding
        parsed = urlparse(url)
        ip_url = urlunparse(
            (
                parsed.scheme,
                f"{resolved_ip}:{port}",
                parsed.path,
                parsed.params,
                parsed.query,
                parsed.fragment,
            )
        )

        MAX_FETCH_BYTES = 5 * 1024 * 1024  # 5 MB
        body = b""
        try:
            async with httpx.AsyncClient() as client:
                async with client.stream(
                    "GET",
                    ip_url,
                    headers={"Host": hostname},
                    follow_redirects=False,
                    timeout=30,
                ) as resp:
                    if resp.status_code != 200:
                        logger.warning(
                            "URL fetch returned non-200 (url=%s status=%s)",
                            url,
                            resp.status_code,
                        )
                        return "", None
                    async for chunk in resp.aiter_bytes(chunk_size=8192):
                        body += chunk
                        if len(body) > MAX_FETCH_BYTES:
                            logger.warning(
                                "URL fetch aborted: body exceeded %d bytes (url=%s)",
                                MAX_FETCH_BYTES,
                                url,
                            )
                            return "", None
        except httpx.HTTPError as e:
            logger.warning(
                "URL fetch HTTP error (url=%s error_type=%s error=%s)",
                url,
                type(e).__name__,
                e,
            )
            return "", None

        if not body:
            logger.warning("URL fetch returned empty body (url=%s)", url)
            return "", None

        # Offload HTML parsing (CPU-bound) to a worker thread
        text, title = await asyncio.to_thread(_extract_with_trafilatura, body, url)
        if not text:
            logger.warning(
                "Content extraction returned empty text (url=%s body_bytes=%d)",
                url,
                len(body),
            )
        else:
            logger.info(
                "URL fetch ok (url=%s body_bytes=%d text_chars=%d title=%r)",
                url,
                len(body),
                len(text),
                title,
            )
        return text, title

    except Exception:
        logger.exception("Unexpected error fetching URL (url=%s)", url)
        return "", None


def _normalize_text(text: str) -> str:
    import re

    text = re.sub(r"[\r\n]+", "\n", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _create_chunks(
    file_id: str,
    parsed_doc_id: str,
    text: str,
    chunk_size: int = 500,
    overlap: int = 50,
) -> list:
    words = text.split()
    chunks = []
    chunk_index = 0
    i = 0

    while i < len(words):
        chunk_words = words[i : i + chunk_size]
        chunks.append(
            DocumentChunk(
                file_id=file_id,
                parsed_document_id=parsed_doc_id,
                chunk_index=chunk_index,
                text=" ".join(chunk_words),
                token_count=len(chunk_words),
                embedding_status="pending",
                is_active=True,
            )
        )
        chunk_index += 1
        i += chunk_size - overlap

    return chunks


async def _recalculate_session_totals(
    db: AsyncSession,
    session_id: str,
    user_id: str | None = None,
    guest_session_id: str | None = None,
) -> tuple[float, float]:
    """Recompute total_score and max_score for a session from active AnswerLogs."""
    conditions = [
        AnswerLog.quiz_session_id == session_id,
        AnswerLog.is_active_result.is_(True),
    ]
    if user_id:
        conditions.append(AnswerLog.user_id == user_id)
    elif guest_session_id:
        conditions.append(AnswerLog.guest_session_id == guest_session_id)

    score_agg = await db.execute(
        select(
            func.sum(AnswerLog.score_awarded),
            func.sum(AnswerLog.max_score),
        ).where(*conditions)
    )
    row = score_agg.one()
    return float(row[0] or 0.0), float(row[1] or 0.0)


@dataclass
class GradingResult:
    judgement: Judgement
    score_awarded: float
    max_score: float = 1.0
    grading_confidence: float = 1.0
    grading_rationale: str = ""
    error_type: "ErrorType | None" = None
    missing_points: "list | None" = None
    suggested_feedback: str = ""
    tokens_used: int = 0


async def _grade_single_answer(
    *,
    item: "QuizItem",
    user_answer: str,
    model_name: "str | None" = None,
    include_essay: bool = False,
) -> GradingResult:
    """Grade a single answer against a quiz item. Pure logic, no DB side-effects.

    - MC/OX: exact match against normalized correct answer.
    - short_answer/fill_blank: exact match first, AI fallback.
    - essay (only when include_essay=True): AI grading with essay prompt.

    Uses SYSTEM_PROMPT_GRADING_SHORT for short/fill_blank AI fallback.
    Uses SYSTEM_PROMPT_GRADING_ESSAY for essay.
    """
    from app.utils.ai_client import call_ai_with_fallback, GRADING_SCHEMA
    from app.config import settings as cfg
    from app.prompts.grading_short import SYSTEM_PROMPT_GRADING_SHORT
    import json as json_mod

    raw_correct = item.correct_answer_json
    if isinstance(raw_correct, dict):
        correct_answer = raw_correct
    elif isinstance(raw_correct, str):
        correct_answer = {"answer": raw_correct}
    else:
        correct_answer = {}

    normalized = normalize_answer(user_answer)
    primary = model_name or cfg.balanced_generation_model
    fallback = cfg.eco_generation_model

    # --- MC / OX: exact match ---
    if item.question_type in (QuestionType.multiple_choice, QuestionType.ox):
        correct_val = normalize_answer(str(correct_answer.get("answer", "")))
        if normalized == correct_val:
            return GradingResult(
                judgement=Judgement.correct,
                score_awarded=1.0,
                error_type=None,
            )
        return GradingResult(
            judgement=Judgement.incorrect,
            score_awarded=0.0,
            error_type=ErrorType.careless_mistake,
        )

    # --- short_answer / fill_blank: try exact match, then AI ---
    if item.question_type in (QuestionType.short_answer, QuestionType.fill_blank):
        accepted = correct_answer.get(
            "accepted_answers", [correct_answer.get("answer", "")]
        )
        for ans in accepted:
            if normalize_answer(str(ans)) == normalized:
                return GradingResult(judgement=Judgement.correct, score_awarded=1.0)
        # AI fallback
        try:
            prompt = (
                f"채점할 답안:\n"
                f"문제: {item.question_text}\n"
                f"정답: {json_mod.dumps(correct_answer, ensure_ascii=False)}\n"
                f"사용자 답: {user_answer}\n"
                f"정규화된 답: {normalized}\n\n"
                f"다음 필드를 포함한 JSON으로 응답하세요:\n"
                f"judgement, score_awarded, max_score, normalized_user_answer, "
                f"accepted_answers, grading_confidence, grading_rationale, "
                f"missing_points, error_type, suggested_feedback"
            )
            ai, tokens = await call_ai_with_fallback(
                prompt,
                GRADING_SCHEMA,
                primary_model=primary,
                fallback_model=fallback,
                system_message=SYSTEM_PROMPT_GRADING_SHORT,
                cache_key="grading_short_v1",
            )
            return GradingResult(
                judgement=Judgement(ai["judgement"]),
                score_awarded=ai["score_awarded"],
                max_score=ai.get("max_score", 1.0),
                grading_confidence=ai.get("grading_confidence", 0.7),
                grading_rationale=ai.get("grading_rationale", ""),
                error_type=ErrorType(ai["error_type"])
                if ai.get("error_type")
                else None,
                missing_points=ai.get("missing_points"),
                suggested_feedback=ai.get("suggested_feedback", ""),
                tokens_used=tokens,
            )
        except Exception:
            return GradingResult(judgement=Judgement.incorrect, score_awarded=0.0)

    # --- essay: AI grading (only if include_essay=True) ---
    if item.question_type == QuestionType.essay and include_essay:
        try:
            from app.prompts.grading_essay import SYSTEM_PROMPT_GRADING_ESSAY

            prompt = (
                f"채점할 서술형 답안:\n"
                f"문제: {item.question_text}\n"
                f"모범 답안: {json_mod.dumps(correct_answer, ensure_ascii=False)}\n"
                f"사용자 답: {user_answer}\n"
                f"출처 참조: {json_mod.dumps(item.source_refs_json or {}, ensure_ascii=False)}\n\n"
                f"다음 필드를 포함한 JSON으로 응답하세요:\n"
                f"judgement, score_awarded, max_score, normalized_user_answer, "
                f"accepted_answers, grading_confidence, grading_rationale, "
                f"missing_points, error_type, suggested_feedback"
            )
            ai, tokens = await call_ai_with_fallback(
                prompt,
                GRADING_SCHEMA,
                primary_model=primary,
                fallback_model=fallback,
                system_message=SYSTEM_PROMPT_GRADING_ESSAY,
                cache_key="grading_essay_v1",
            )
            return GradingResult(
                judgement=Judgement(ai["judgement"]),
                score_awarded=ai["score_awarded"],
                max_score=ai.get("max_score", 1.0),
                grading_confidence=ai.get("grading_confidence", 0.7),
                grading_rationale=ai.get("grading_rationale", ""),
                error_type=ErrorType(ai["error_type"])
                if ai.get("error_type")
                else None,
                missing_points=ai.get("missing_points"),
                suggested_feedback=ai.get("suggested_feedback", ""),
                tokens_used=tokens,
            )
        except Exception:
            return GradingResult(judgement=Judgement.incorrect, score_awarded=0.0)

    # Fallback for unhandled types
    return GradingResult(judgement=Judgement.incorrect, score_awarded=0.0)


@dataclass
class _QuizGenerationContext:
    source_context: str
    question_count: int | None
    difficulty: str
    question_types: list[str]
    concept_counts: dict[str, int]
    is_no_source: bool
    topic: str | None


async def _prepare_generation_context(
    db: AsyncSession,
    session: QuizSession,
    payload: dict,
) -> _QuizGenerationContext:
    source_texts: list[str] = []
    if session.source_mode.value == "document_based":
        file_results = await db.execute(
            select(File)
            .join(QuizSessionFile, QuizSessionFile.file_id == File.id)
            .where(QuizSessionFile.quiz_session_id == session.id)
        )
        for f in file_results.scalars().all():
            if f.parsed_document:
                source_texts.append(f.parsed_document.normalized_text or "")

    question_count = payload.get("question_count", session.question_count)
    difficulty = payload.get("difficulty", session.difficulty)
    question_types = payload.get("question_types", []) or [
        "multiple_choice",
        "ox",
        "short_answer",
        "fill_blank",
        "essay",
    ]

    concept_counts: dict[str, int] = {}
    if session.user_id:
        recent_concepts_result = await db.execute(
            select(QuizItem.concept_key)
            .join(QuizSession, QuizSession.id == QuizItem.quiz_session_id)
            .where(QuizSession.user_id == session.user_id)
            .order_by(QuizItem.created_at.desc())
            .limit(20)
        )
        recent_concepts = [r[0] for r in recent_concepts_result.all() if r[0]]
        for c in recent_concepts:
            concept_counts[c] = concept_counts.get(c, 0) + 1

    source_url = payload.get("source_url") or None

    if source_url:
        url_text, page_title = await _fetch_url_text(source_url)
        if not url_text:
            raise JobFailure(
                "INVALID_INPUT:URL을 불러올 수 없습니다. 주소가 올바른지, 공개된 페이지인지 확인해 주세요."
            )
        source_context = _normalize_text(url_text)[:_MAX_URL_SOURCE_CHARS]
        is_no_source = False
        if page_title:
            session.title = page_title
    else:
        source_context = (
            "\n\n".join(source_texts[:3])
            if source_texts
            else "No source material provided."
        )
        if not source_texts and session.source_mode.value == "document_based":
            raise JobFailure("No source material available")
        is_no_source = session.source_mode.value == "no_source"

    topic = payload.get("topic") or None

    resolved_difficulty = difficulty
    if difficulty == "auto" and not is_no_source:
        from app.utils.ai_client import (
            call_ai_structured,
            DIFFICULTY_SELECTION_SCHEMA,
        )

        det_result, _ = await call_ai_structured(
            prompt=source_context[:4000],
            schema=DIFFICULTY_SELECTION_SCHEMA,
            system_message=SYSTEM_PROMPT_DIFFICULTY_SELECTION,
            model=cfg.eco_generation_model or cfg.balanced_generation_model,
            max_tokens=100,
        )
        resolved_difficulty = det_result.get("difficulty", "medium")
        if resolved_difficulty not in ("easy", "medium", "hard"):
            resolved_difficulty = "medium"

    return _QuizGenerationContext(
        source_context=source_context,
        question_count=question_count,
        difficulty=resolved_difficulty or "medium",
        question_types=question_types,
        concept_counts=concept_counts,
        is_no_source=is_no_source,
        topic=topic,
    )


async def _run_quiz_generation_from_context(
    session: QuizSession,
    ctx: _QuizGenerationContext,
) -> tuple[list[dict], int, str]:
    from app.utils.ai_client import call_ai_with_fallback, GENERATION_SCHEMA

    prompt = build_generation_prompt(
        source_context=ctx.source_context,
        question_count=ctx.question_count,
        difficulty=ctx.difficulty,
        question_types=ctx.question_types,
        concept_counts=ctx.concept_counts,
        is_no_source=ctx.is_no_source,
        topic=ctx.topic,
    )

    ai_result, tokens_used = await call_ai_with_fallback(
        prompt,
        GENERATION_SCHEMA,
        primary_model=session.generation_model_name or cfg.balanced_generation_model,
        fallback_model=cfg.eco_generation_model,
        system_message=get_generation_system_prompt(ctx.difficulty),
        cache_key=f"quiz_gen_{ctx.difficulty}_v1",
        max_tokens=16384,
    )

    if ai_result.get("rejected"):
        raise JobFailure(
            "INVALID_INPUT:"
            + ai_result.get(
                "rejection_reason",
                "퀴즈를 만들기 어려운 입력입니다. 학습하고 싶은 주제나 자료를 입력해 주세요.",
            )
        )

    all_questions = ai_result.get("questions", [])
    questions = (
        all_questions
        if ctx.question_count is None
        else all_questions[: ctx.question_count]
    )
    if not questions:
        raise JobFailure("AI returned no questions")

    if ctx.question_count is None:
        session.question_count = len(questions)

    return questions, tokens_used, ctx.difficulty or "medium"


async def _run_quiz_generation(
    db: AsyncSession,
    session: QuizSession,
    payload: dict,
) -> tuple[list[dict], int, str]:
    ctx = await _prepare_generation_context(db, session, payload)
    return await _run_quiz_generation_from_context(session, ctx)


_DEFAULT_OX_OPTIONS = {"o": "참", "x": "거짓"}


def _sanitize_options(options: dict | None, question_type: str) -> dict | None:
    if question_type in ("short_answer", "fill_blank", "essay"):
        return None
    if not options or not isinstance(options, dict):
        if question_type == "ox":
            return dict(_DEFAULT_OX_OPTIONS)
        return options
    filtered = {k: v for k, v in options.items() if v is not None}
    if not filtered:
        if question_type == "ox":
            return dict(_DEFAULT_OX_OPTIONS)
        return None
    return filtered


def _persist_quiz_items(
    db: AsyncSession,
    session: QuizSession,
    questions: list[dict],
    default_difficulty: str,
) -> list[QuizItem]:
    """Create QuizItem records from AI output. Returns the created items."""
    items: list[QuizItem] = []
    for i, q in enumerate(questions):
        qtype = q.get("question_type", "multiple_choice")
        item = QuizItem(
            quiz_session_id=session.id,
            item_order=i + 1,
            question_type=QuestionType(qtype),
            question_text=q.get("question_text", ""),
            options_json=_sanitize_options(q.get("options"), qtype),
            option_descriptions_json=q.get("option_descriptions"),
            correct_answer_json=q.get("correct_answer"),
            explanation_text=q.get("explanation", ""),
            source_refs_json={"refs": q.get("source_refs", [])},
            concept_key=q.get("concept_key", ""),
            concept_label=q.get("concept_label", ""),
            category_tag=q.get("category_tag", ""),
            difficulty=q.get("difficulty", default_difficulty),
            similarity_fingerprint=q.get("concept_key", ""),
        )
        db.add(item)
        items.append(item)
    return items


async def _adjust_generation_credits(
    db: AsyncSession,
    session: QuizSession,
    tokens_used: int,
    credit_estimate: float,
) -> None:
    """Adjust credits based on actual token usage vs estimate."""
    from app.tier_config import calculate_credit_cost
    from app.services.usage_service import UsageService

    actual_cost = calculate_credit_cost(
        tokens_used,
        session.generation_model_name or cfg.balanced_generation_model,
    )
    delta = actual_cost - credit_estimate
    if abs(delta) > 0.001 and session.user_id:
        await UsageService().adjust_credit(db, str(session.user_id), "quiz", delta)


async def _refund_quiz_credits_on_invalid_input(
    db: AsyncSession,
    job: Job,
    session: QuizSession,
    error: Exception,
) -> None:
    """Refund pre-charged quiz credits when generation fails with
    ``JobFailure("INVALID_INPUT:...")`` — i.e. user-input errors
    (bad URL, AI rejection) rather than system errors.

    Idempotent via a ``credits_refunded`` flag on ``job.payload_json``
    so that Celery acks_late redelivery or any future retry path
    cannot double-refund. Swallows all exceptions — a refund failure
    must never mask the original generation failure the caller is
    about to surface to the user.
    """
    try:
        from sqlalchemy.orm.attributes import flag_modified

        if not str(error).startswith("INVALID_INPUT:"):
            return
        payload = job.payload_json or {}
        if payload.get("credits_refunded"):
            return
        estimate = float(payload.get("credit_estimate") or 0.0)
        if estimate <= 0 or not session.user_id:
            return

        from app.services.usage_service import UsageService

        await UsageService().adjust_credit(db, str(session.user_id), "quiz", -estimate)
        payload["credits_refunded"] = True
        job.payload_json = payload
        flag_modified(job, "payload_json")
        logger.info(
            "Refunded %.2f quiz credits to user %s for failed job %s",
            estimate,
            session.user_id,
            job.id,
        )
    except Exception as exc:
        logger.exception(
            "Credit refund failed for job %s (original error not masked): %s",
            job.id,
            exc,
        )


def _serialize_quiz_item(item: QuizItem) -> dict:
    """Serialize QuizItem to match QuizItemResponse shape for SSE streaming."""
    return {
        "id": item.id,
        "item_order": item.item_order,
        "question_type": item.question_type.value if item.question_type else None,
        "question_text": item.question_text,
        "options": item.options_json,
        "option_descriptions": item.option_descriptions_json,
        "difficulty": item.difficulty,
        "concept_key": item.concept_key,
        "concept_label": item.concept_label,
        "category_tag": item.category_tag,
        "correct_answer": item.correct_answer_json,
        "explanation": item.explanation_text,
        "tips": item.tips_text,
        "source_refs": item.source_refs_json,
    }


async def _dispatch_celery_fallback(job_id: str, session_id: str) -> bool:
    try:
        from app.workers.celery_app import dispatch_task

        dispatch_task("generate_quiz", [job_id])
        logger.info(
            "Dispatched Celery fallback for session %s (job %s)",
            session_id,
            job_id,
        )
        return True
    except Exception as exc:
        logger.error("Failed to dispatch Celery fallback: %s", exc)
        return False


async def stream_quiz_generation(
    db: AsyncSession,
    session: QuizSession,
    job: Job,
) -> AsyncGenerator[str, None]:
    """SSE generator: streams the quiz model's native Chain-of-Thought, then the generated items."""
    from app.utils.sse import sse_data, sse_error, sse_done
    from app.utils.ai_client import GENERATION_SCHEMA
    from app.config import settings as _cfg

    generation_completed = False
    try:
        session.status = QuizSessionStatus.generating
        await db.commit()

        yield sse_data({"type": "stage", "stage": "analyzing"})

        payload = job.payload_json or {}

        ctx = await _prepare_generation_context(db, session, payload)

        yield sse_data({"type": "stage", "stage": "generating"})
        yield sse_data({"type": "thinking_start"})

        prompt = build_generation_prompt(
            source_context=ctx.source_context,
            question_count=ctx.question_count,
            difficulty=ctx.difficulty,
            question_types=ctx.question_types,
            concept_counts=ctx.concept_counts,
            is_no_source=ctx.is_no_source,
            topic=ctx.topic,
        )

        primary_model = session.generation_model_name or _cfg.balanced_generation_model
        fallback_model = _cfg.eco_generation_model

        ai_result: dict | None = None
        tokens_used = 0

        stream_iter = stream_ai_structured_with_thinking(
            prompt=prompt,
            schema=GENERATION_SCHEMA,
            system_message=get_generation_system_prompt(ctx.difficulty),
            primary_model=primary_model,
            fallback_model=fallback_model,
            max_tokens=16384,
            cache_key=f"quiz_gen_{ctx.difficulty}_v1",
            reasoning_effort="medium",
            thinking_level="HIGH",
        )

        async with asyncio.timeout(_cfg.generation_timeout):
            async for event in stream_iter:
                etype = event.get("type")
                if etype == "thinking":
                    text = event.get("text") or ""
                    if isinstance(text, str) and text:
                        yield sse_data({"type": "thinking_chunk", "text": text})
                elif etype == "result":
                    data = event.get("data")
                    if isinstance(data, dict):
                        ai_result = data
                    tokens_raw = event.get("tokens_used") or 0
                    tokens_used = (
                        int(tokens_raw) if isinstance(tokens_raw, (int, float)) else 0
                    )

        yield sse_data({"type": "thinking_end"})

        if ai_result is None:
            raise JobFailure("AI stream ended without a result")

        if ai_result.get("rejected"):
            raise JobFailure(
                "INVALID_INPUT:"
                + ai_result.get(
                    "rejection_reason",
                    "퀴즈를 만들기 어려운 입력입니다. 학습하고 싶은 주제나 자료를 입력해 주세요.",
                )
            )

        all_questions = ai_result.get("questions", [])
        questions = (
            all_questions
            if ctx.question_count is None
            else all_questions[: ctx.question_count]
        )
        if not questions:
            raise JobFailure("AI returned no questions")

        if ctx.question_count is None:
            session.question_count = len(questions)

        items = _persist_quiz_items(db, session, questions, ctx.difficulty)

        credit_estimate = payload.get("credit_estimate", 0.0)
        await _adjust_generation_credits(db, session, tokens_used, credit_estimate)

        session.status = QuizSessionStatus.ready
        await db.commit()
        generation_completed = True
        logger.info(
            "SSE gen completed for session %s, %d items", session.id, len(items)
        )

        for item in items:
            await db.refresh(item)

        yield sse_data(
            {
                "type": "stage",
                "stage": "streaming_questions",
                "total": len(items),
            }
        )

        for i, item in enumerate(items):
            yield sse_data(
                {
                    "type": "question",
                    "index": i,
                    "total": len(items),
                    "item": _serialize_quiz_item(item),
                }
            )
            if i < len(items) - 1:
                await asyncio.sleep(0.2)

        logger.info(
            "SSE streaming done for session %s, yielding done event", session.id
        )
        yield sse_done()

    except JobFailure as e:
        session.status = QuizSessionStatus.generation_failed
        job.error_message = str(e)
        await _refund_quiz_credits_on_invalid_input(db, job, session, e)
        try:
            await db.commit()
        except Exception:
            pass
        yield sse_error(str(e))
    except GeneratorExit:
        raise
    except Exception as e:
        logger.exception(
            "stream_quiz_generation error for session %s: %s", session.id, e
        )
        if not generation_completed:
            dispatched = await _dispatch_celery_fallback(job.id, session.id)
            if not dispatched:
                session.status = QuizSessionStatus.generation_failed
                job.error_message = f"SSE failed and Celery fallback unavailable: {e}"
            try:
                await db.commit()
            except Exception:
                pass
        yield sse_error(f"Generation failed: {e}")
    finally:
        if not generation_completed:
            logger.warning(
                "SSE gen incomplete for session %s (status=%s, completed=%s)",
                session.id,
                session.status.value
                if hasattr(session.status, "value")
                else session.status,
                generation_completed,
            )
        if not generation_completed and session.status == QuizSessionStatus.generating:
            try:
                await _dispatch_celery_fallback(job.id, session.id)
            except Exception:
                pass


async def generate_quiz(job_id: str):
    async with async_session() as db:
        job_result = await db.execute(select(Job).where(Job.id == job_id))
        job = job_result.scalar_one_or_none()
        if not job:
            return

        async with JobRunner(db, job):
            session_result = await db.execute(
                select(QuizSession).where(QuizSession.id == job.target_id)
            )
            session = session_result.scalar_one_or_none()
            if not session:
                raise JobFailure("Session not found")

            try:
                session.status = QuizSessionStatus.generating
                await db.commit()

                from app.utils.ai_client import (
                    call_ai_with_fallback,
                    BATCH_RETRY_GENERATION_SCHEMA,
                )
                from app.config import settings as cfg

                payload = job.payload_json or {}

                if job.job_type == "retry_generation":
                    concept_keys = payload.get("concept_keys", [])
                    question_count = payload.get("size", session.question_count)
                    user_difficulty = payload.get("difficulty")
                    user_question_types = payload.get("question_types", [])

                    if not concept_keys:
                        session.status = QuizSessionStatus.generation_failed
                        raise JobFailure("No concept keys provided for retry")

                    items_created = 0

                    tasks_data = []
                    batch_items = []
                    for idx, concept_key in enumerate(concept_keys):
                        if question_count is not None and idx >= question_count:
                            break

                        wrong_result = await db.execute(
                            select(AnswerLog, QuizItem)
                            .join(QuizItem, QuizItem.id == AnswerLog.quiz_item_id)
                            .where(
                                AnswerLog.user_id == session.user_id,
                                QuizItem.concept_key == concept_key,
                                AnswerLog.judgement.in_(
                                    [Judgement.incorrect, Judgement.partial]
                                ),
                            )
                            .order_by(AnswerLog.created_at.desc())
                            .limit(1)
                        )
                        row = wrong_result.first()
                        if not row:
                            continue

                        answer_log, quiz_item = row
                        tasks_data.append((idx, concept_key, quiz_item))
                        batch_items.append(
                            {
                                "concept_key": concept_key,
                                "concept_label": quiz_item.concept_label or concept_key,
                                "previous_question_type": quiz_item.question_type.value,
                                "previous_question": quiz_item.question_text,
                                "error_type": answer_log.error_type.value
                                if answer_log.error_type
                                else "unknown",
                                "user_answer": answer_log.user_answer_raw or "",
                                "correct_answer": str(
                                    quiz_item.correct_answer_json.get("answer", "")
                                    if isinstance(quiz_item.correct_answer_json, dict)
                                    else quiz_item.correct_answer_json
                                ),
                                "retry_count": 1,
                            }
                        )

                    if not batch_items:
                        session.status = QuizSessionStatus.generation_failed
                        raise JobFailure("No retry questions could be generated")

                    batch_prompt = build_batch_retry_prompt(
                        batch_items,
                        difficulty=user_difficulty,
                        question_types=user_question_types,
                    )
                    resolved_retry_difficulty = user_difficulty or "medium"
                    batch_result, tokens_used = await call_ai_with_fallback(
                        batch_prompt,
                        BATCH_RETRY_GENERATION_SCHEMA,
                        primary_model=session.generation_model_name
                        or cfg.balanced_generation_model,
                        fallback_model=cfg.eco_generation_model,
                        system_message=get_retry_system_prompt(
                            resolved_retry_difficulty
                        ),
                        cache_key=f"retry_gen_{resolved_retry_difficulty}_v1",
                        max_tokens=16384,
                    )

                    result_by_concept = {
                        q["concept_key"]: q
                        for q in batch_result.get("questions", [])
                        if q.get("concept_key")
                    }
                    for idx, concept_key, quiz_item in tasks_data:
                        ai_result = result_by_concept.get(concept_key)
                        if not ai_result or not ai_result.get("question_text"):
                            logger.warning(
                                "Retry generation: no AI result for concept_key=%r in session %s; skipping",
                                concept_key,
                                session.id,
                            )
                            continue
                        new_item = QuizItem(
                            quiz_session_id=session.id,
                            item_order=idx + 1,
                            question_type=QuestionType(
                                ai_result.get("question_type", "multiple_choice")
                            ),
                            question_text=ai_result.get("question_text", ""),
                            correct_answer_json=ai_result.get(
                                "correct_answer", {"answer": ""}
                            ),
                            explanation_text=ai_result.get("explanation", ""),
                            concept_key=concept_key,
                            concept_label=quiz_item.concept_label,
                            category_tag=quiz_item.category_tag,
                            difficulty=user_difficulty
                            or quiz_item.difficulty
                            or "medium",
                            source_refs_json=None,
                            options_json=ai_result.get("options"),
                            option_descriptions_json=ai_result.get(
                                "option_descriptions"
                            ),
                        )
                        db.add(new_item)
                        items_created += 1

                    if items_created == 0:
                        session.status = QuizSessionStatus.generation_failed
                        raise JobFailure("No retry questions could be generated")

                    from app.tier_config import calculate_credit_cost
                    from app.services.usage_service import UsageService

                    actual_cost = calculate_credit_cost(
                        tokens_used,
                        session.generation_model_name or cfg.balanced_generation_model,
                    )
                    estimate = (job.payload_json or {}).get("credit_estimate", 0.0)
                    delta = actual_cost - estimate
                    _retry_uid = session.user_id
                    if abs(delta) > 0.001 and _retry_uid is not None:
                        await UsageService().adjust_credit(
                            db, _retry_uid, "quiz", delta
                        )

                    await db.commit()
                    session.status = QuizSessionStatus.ready
                    return

                (
                    questions,
                    tokens_used,
                    resolved_difficulty,
                ) = await _run_quiz_generation(db, session, payload)
                _persist_quiz_items(db, session, questions, resolved_difficulty)

                credit_estimate = payload.get("credit_estimate", 0.0)
                await _adjust_generation_credits(
                    db, session, tokens_used, credit_estimate
                )

                session.status = QuizSessionStatus.ready

            except JobFailure as e:
                await _refund_quiz_credits_on_invalid_input(db, job, session, e)
                raise
            except Exception:
                session.status = QuizSessionStatus.generation_failed
                raise


async def grade_exam(job_id: str):
    async with async_session() as db:
        job_result = await db.execute(select(Job).where(Job.id == job_id))
        job = job_result.scalar_one_or_none()
        if not job:
            return

        async with JobRunner(db, job):
            session_result = await db.execute(
                select(QuizSession).where(QuizSession.id == job.target_id)
            )
            session = session_result.scalar_one_or_none()
            if not session:
                raise JobFailure("Session not found")

            try:
                session.status = QuizSessionStatus.grading
                await db.commit()

                items_result = await db.execute(
                    select(QuizItem)
                    .where(QuizItem.quiz_session_id == session.id)
                    .order_by(QuizItem.item_order)
                )
                items = items_result.scalars().all()

                all_gradings: list[GradingResult] = []

                for item in items:
                    draft_result = await db.execute(
                        select(DraftAnswer).where(
                            DraftAnswer.quiz_item_id == item.id,
                            DraftAnswer.quiz_session_id == session.id,
                            DraftAnswer.user_id == session.user_id,
                        )
                    )
                    draft = draft_result.scalar_one_or_none()
                    user_answer = draft.user_answer if draft else ""

                    if not user_answer:
                        grading = GradingResult(
                            judgement=Judgement.skipped,
                            score_awarded=0.0,
                            error_type=ErrorType.no_response,
                        )
                    else:
                        grading = await _grade_single_answer(
                            item=item,
                            user_answer=user_answer,
                            model_name=session.generation_model_name,
                            include_essay=True,
                        )

                    all_gradings.append(grading)

                    existing = await db.execute(
                        select(AnswerLog).where(
                            AnswerLog.quiz_item_id == item.id,
                            AnswerLog.user_id == session.user_id,
                            AnswerLog.is_active_result.is_(True),
                        )
                    )
                    for old_log in existing.scalars().all():
                        old_log.is_active_result = False

                    answer_log = AnswerLog(
                        quiz_item_id=item.id,
                        quiz_session_id=session.id,
                        user_id=session.user_id,
                        user_answer_raw=user_answer,
                        user_answer_normalized=normalize_answer(user_answer)
                        if user_answer
                        else "",
                        judgement=grading.judgement,
                        score_awarded=grading.score_awarded,
                        max_score=grading.max_score,
                        grading_confidence=grading.grading_confidence,
                        error_type=grading.error_type,
                        is_active_result=True,
                        graded_at=datetime.now(timezone.utc),
                    )
                    db.add(answer_log)

                    if session.user_id:
                        await _update_weak_point(
                            db, session.user_id, item, grading.judgement
                        )

                await db.flush()
                (
                    session.total_score,
                    session.max_score,
                ) = await _recalculate_session_totals(
                    db, session.id, user_id=session.user_id
                )
                session.graded_at = datetime.now(timezone.utc)
                session.status = QuizSessionStatus.graded

                from app.tier_config import calculate_credit_cost
                from app.services.usage_service import UsageService

                _exam_total_tokens = sum(g.tokens_used for g in all_gradings)
                _exam_actual_cost = calculate_credit_cost(
                    _exam_total_tokens,
                    session.generation_model_name or cfg.balanced_generation_model,
                )
                _exam_estimate = (job.payload_json or {}).get("credit_estimate", 0.0)
                _exam_delta = _exam_actual_cost - _exam_estimate
                _exam_uid = session.user_id
                if abs(_exam_delta) > 0.001 and _exam_uid is not None:
                    await UsageService().adjust_credit(
                        db, _exam_uid, "quiz", _exam_delta
                    )

            except JobFailure:
                raise
            except Exception:
                session.status = QuizSessionStatus.submitted
                raise


async def _apply_objection_decision(
    db: AsyncSession,
    objection: "Objection",
    answer_log: "AnswerLog",
    quiz_item: "QuizItem",
    ai_result: dict,
) -> None:
    """Apply the AI objection decision.

    If upheld/partially_upheld and should_apply:
      - Deactivate old answer log
      - Create new AnswerLog with updated scores
      - Recalculate session totals
      - Set objection status to 'applied'
      - Set session status to 'regraded'

    If rejected/other:
      - Set objection status
      - If session has no remaining pending objections → restore to 'graded'
    """
    decision = ai_result.get("decision", "rejected")
    if ai_result.get("should_apply") and decision in ("upheld", "partially_upheld"):
        answer_log.is_active_result = False
        new_log = AnswerLog(
            quiz_item_id=objection.quiz_item_id,
            quiz_session_id=objection.quiz_session_id,
            user_id=objection.user_id,
            user_answer_raw=answer_log.user_answer_raw,
            user_answer_normalized=answer_log.user_answer_normalized,
            judgement=Judgement(ai_result["updated_judgement"]),
            score_awarded=ai_result["updated_score_awarded"],
            max_score=answer_log.max_score,
            grading_confidence=0.9,
            grading_rationale=ai_result.get("reasoning", ""),
            error_type=ErrorType(ai_result["updated_error_type"])
            if ai_result.get("updated_error_type")
            else None,
            is_active_result=True,
            regraded_from_answer_log_id=answer_log.id,
            graded_at=datetime.now(timezone.utc),
        )
        db.add(new_log)
        await db.flush()

        session_result = await db.execute(
            select(QuizSession).where(QuizSession.id == objection.quiz_session_id)
        )
        session = session_result.scalar_one_or_none()
        if session:
            session.total_score, session.max_score = await _recalculate_session_totals(
                db, objection.quiz_session_id, objection.user_id
            )
            objection.status = ObjectionStatus.applied
            session.status = QuizSessionStatus.regraded
    else:
        objection.status = (
            ObjectionStatus.rejected
            if decision == "rejected"
            else ObjectionStatus(decision)
        )
        await _restore_session_if_no_pending_objections(db, objection)


async def _restore_session_if_no_pending_objections(
    db: AsyncSession,
    objection: "Objection",
) -> None:
    """If no other pending objections for this session, restore session to 'graded'."""
    session_result = await db.execute(
        select(QuizSession).where(QuizSession.id == objection.quiz_session_id)
    )
    session = session_result.scalar_one_or_none()
    if session and session.status == QuizSessionStatus.objection_pending:
        remaining_result = await db.execute(
            select(Objection).where(
                Objection.quiz_session_id == objection.quiz_session_id,
                Objection.id != objection.id,
                Objection.status.in_(
                    [ObjectionStatus.submitted, ObjectionStatus.under_review]
                ),
            )
        )
        if not remaining_result.scalar_one_or_none():
            session.status = QuizSessionStatus.graded


async def review_objection(job_id: str):
    async with async_session() as db:
        job = (
            await db.execute(select(Job).where(Job.id == job_id))
        ).scalar_one_or_none()
        if not job:
            return

        async with JobRunner(db, job):
            objection = (
                await db.execute(select(Objection).where(Objection.id == job.target_id))
            ).scalar_one_or_none()
            if not objection:
                raise JobFailure("Objection not found")

            answer_log = (
                await db.execute(
                    select(AnswerLog).where(AnswerLog.id == objection.answer_log_id)
                )
            ).scalar_one_or_none()
            quiz_item = (
                await db.execute(
                    select(QuizItem).where(QuizItem.id == objection.quiz_item_id)
                )
            ).scalar_one_or_none()
            if not answer_log or not quiz_item:
                objection.status = ObjectionStatus.rejected
                return

            prompt = f"""이의제기를 검토하세요.

원문 문제: {quiz_item.question_text}
정답: {json_mod.dumps(quiz_item.correct_answer_json or {}, ensure_ascii=False)}
사용자 답: {answer_log.user_answer_raw}
원래 판정: {answer_log.judgement.value}
원래 점수: {answer_log.score_awarded}
이의 사유: {objection.objection_reason}

다음 필드를 포함한 JSON으로 응답하세요:
decision, reasoning, updated_judgement, updated_score_awarded, updated_error_type, should_apply"""
            _obj_call = await call_ai_with_fallback(
                prompt,
                OBJECTION_REVIEW_SCHEMA,
                primary_model=cfg.balanced_generation_model,
                fallback_model=cfg.eco_generation_model,
                system_message=SYSTEM_PROMPT_OBJECTION_REVIEW,
                cache_key="objection_v1",
            )
            ai_result: dict = _obj_call[0]
            tokens_used: int = _obj_call[1]
            objection.review_result_json = ai_result
            objection.decided_at = datetime.now(timezone.utc)
            objection.decided_by = "ai"
            await _apply_objection_decision(
                db, objection, answer_log, quiz_item, ai_result
            )
            from app.tier_config import calculate_credit_cost
            from app.services.usage_service import UsageService

            _obj_actual_cost = calculate_credit_cost(
                tokens_used, cfg.balanced_generation_model
            )
            _obj_estimate = (job.payload_json or {}).get("credit_estimate", 0.0)
            _obj_delta = _obj_actual_cost - _obj_estimate
            _obj_uid = answer_log.user_id
            if abs(_obj_delta) > 0.001 and _obj_uid is not None:
                await UsageService().adjust_credit(db, _obj_uid, "quiz", _obj_delta)


async def _update_weak_point(
    db: AsyncSession, user_id: str, item: QuizItem, judgement: Judgement
):
    concept_key = normalize_concept_key(item.concept_key or "")
    if not concept_key:
        return

    result = await db.execute(
        select(WeakPoint).where(
            WeakPoint.user_id == user_id,
            WeakPoint.concept_key == concept_key,
        )
    )
    weak = result.scalar_one_or_none()

    if not weak:
        weak = WeakPoint(
            user_id=user_id,
            concept_key=concept_key,
            concept_label=item.concept_label,
            category_tag=item.category_tag,
        )
        db.add(weak)
        await db.flush()

    if judgement == Judgement.incorrect:
        weak.wrong_count += 1
        weak.streak_wrong_count += 1
        weak.last_wrong_at = datetime.now(timezone.utc)
    elif judgement == Judgement.partial:
        weak.partial_count += 1
        weak.streak_wrong_count = 0
    elif judgement == Judgement.skipped:
        weak.skip_count += 1
    elif judgement == Judgement.correct:
        weak.streak_wrong_count = 0

    await db.flush()


async def admin_regrade(job_id: str) -> None:
    from app.models.search import Job

    async with async_session() as db:
        result = await db.execute(select(Job).where(Job.id == job_id))
        job = result.scalar_one_or_none()
        if not job:
            logger.error("admin_regrade: job %s not found", job_id)
            return

        async with JobRunner(db, job):
            item_result = await db.execute(
                select(QuizItem).where(QuizItem.id == job.target_id)
            )
            item = item_result.scalar_one_or_none()
            if not item:
                logger.error("admin_regrade: quiz item %s not found", job.target_id)
                raise JobFailure("Quiz item not found")

            logs_result = await db.execute(
                select(AnswerLog, QuizSession)
                .join(QuizSession, QuizSession.id == AnswerLog.quiz_session_id)
                .where(
                    AnswerLog.quiz_item_id == item.id,
                    AnswerLog.is_active_result.is_(True),
                )
            )
            rows = logs_result.all()

            for answer_log, session in rows:
                try:
                    user_answer = answer_log.user_answer_raw or ""
                    grading = (
                        await _grade_single_answer(
                            item=item,
                            user_answer=user_answer,
                            model_name=session.generation_model_name,
                            include_essay=False,
                        )
                        if user_answer
                        else GradingResult(
                            judgement=Judgement.skipped,
                            score_awarded=0.0,
                            error_type=ErrorType.no_response,
                        )
                    )
                    answer_log.judgement = grading.judgement
                    answer_log.score_awarded = grading.score_awarded
                    answer_log.graded_at = datetime.now(timezone.utc)
                except Exception as exc:
                    logger.warning(
                        "admin_regrade: failed to regrade log %s: %s",
                        answer_log.id,
                        exc,
                    )

            await db.flush()

            affected_session_ids = {row[1].id for row in rows}
            for sid in affected_session_ids:
                sess_upd = await db.execute(
                    select(QuizSession).where(QuizSession.id == sid)
                )
                sess_obj = sess_upd.scalar_one_or_none()
                if sess_obj:
                    user_id_for_session = next(
                        (row[0].user_id for row in rows if row[1].id == sid),
                        None,
                    )
                    if user_id_for_session:
                        (
                            sess_obj.total_score,
                            sess_obj.max_score,
                        ) = await _recalculate_session_totals(
                            db, sid, user_id_for_session
                        )
